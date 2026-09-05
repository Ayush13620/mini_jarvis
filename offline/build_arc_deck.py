import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_cyberpunk_deck():
    prs = Presentation()
    # 16:9 Widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Cyberpunk HUD Palette
    BG_DARK = RGBColor(7, 11, 20)          # Ultra Dark Navy #070B14
    CARD_BG = RGBColor(15, 23, 42)         # Translucent Card #0F172A
    CARD_HEADER_BG = RGBColor(24, 37, 66)  # Dark Blue Header #182542
    NEON_CYAN = RGBColor(0, 240, 255)      # Cyber Cyan #00F0FF
    NEON_BLUE = RGBColor(56, 189, 248)     # Electric Blue #38BDF8
    NEON_GREEN = RGBColor(34, 197, 94)     # Success Green #22C55E
    ALERT_RED = RGBColor(255, 45, 85)      # Alert Red #FF2D55
    TEXT_WHITE = RGBColor(248, 250, 252)   # Slate White #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)   # Slate Muted #94A3B8

    def apply_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_top_bar(slide, title_text, tracker_text="PROJECT A.R.K. | AUTONOMOUS RELAY CONTROL & HOLOGRAPHIC SYSTEM"):
        # Top Tracker / Category Line
        tb_track = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.3))
        p_tr = tb_track.text_frame.paragraphs[0]
        p_tr.text = f"//  {tracker_text.upper()}"
        p_tr.font.size = Pt(9.5)
        p_tr.font.bold = True
        p_tr.font.color.rgb = NEON_CYAN
        p_tr.font.name = "Consolas"

        # Main Slide Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.7))
        p_ti = tb_title.text_frame.paragraphs[0]
        p_ti.text = title_text
        p_ti.font.size = Pt(22)
        p_ti.font.bold = True
        p_ti.font.color.rgb = TEXT_WHITE
        p_ti.font.name = "Arial"

        # Top Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = NEON_CYAN
        line.line.fill.background()

    def add_hud_card(slide, left, top, width, height, title="", items=None, border_color=NEON_BLUE, bg_color=CARD_BG):
        # Card Background Container
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        # Card Text Frame
        tb = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.15), Inches(width - 0.36), Inches(height - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        if title:
            p_t = tf.paragraphs[0]
            p_t.text = title.upper()
            p_t.font.size = Pt(14)
            p_t.font.bold = True
            p_t.font.color.rgb = NEON_CYAN
            p_t.font.name = "Arial"
            p_t.space_after = Pt(8)

        if items:
            first = False if title else True
            for item in items:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()

                if isinstance(item, tuple):
                    p.text = f"> {item[0]}: "
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.color.rgb = NEON_BLUE
                    p.font.name = "Arial"

                    run = p.add_run()
                    run.text = item[1]
                    run.font.bold = False
                    run.font.color.rgb = TEXT_WHITE
                    run.font.name = "Arial"
                else:
                    p.text = f"• {item}"
                    p.font.size = Pt(11)
                    p.font.color.rgb = TEXT_WHITE
                    p.font.name = "Arial"
                p.space_after = Pt(6)

    def add_stat_badge(slide, left, top, width, height, big_val, label_text, color=NEON_CYAN):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(width), Inches(height - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = big_val
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.font.name = "Consolas"

        p2 = tf.add_paragraph()
        p2.text = label_text.upper()
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MUTED
        p2.font.name = "Arial"

    # =========================================================================
    # SLIDE 1: Title Slide (Iron Man Arc Reactor Cyberpunk Theme)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_bg(s1)

    # Big Decorative Outer Frame
    frame = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD_BG
    frame.line.color.rgb = NEON_CYAN
    frame.line.width = Pt(2)

    # Title Text Frame
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "// PROJECT A.R.K."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN
    p.font.name = "Consolas"
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "AUTONOMOUS RELAY CONTROL &\nHOLOGRAPHIC KNOWLEDGE SYSTEM"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Arial"
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "An Edge-Cloud AI Voice Assistant with Pepper's Ghost 4-Axis Arc Reactor Hologram & Hardware Switching"
    p.font.size = Pt(14)
    p.font.color.rgb = NEON_BLUE
    p.font.name = "Arial"
    p.space_after = Pt(24)

    # Bottom Badges
    add_stat_badge(s1, 1.2, 5.0, 2.5, 1.1, "< 800ms", "Cloud Latency", NEON_CYAN)
    add_stat_badge(s1, 4.0, 5.0, 2.5, 1.1, "16 kHz", "TCP Streaming", NEON_BLUE)
    add_stat_badge(s1, 6.8, 5.0, 2.5, 1.1, "0.96\" OLED", "Mini Hologram", NEON_GREEN)
    add_stat_badge(s1, 9.6, 5.0, 2.5, 1.1, "5V Relay", "Appliance Control", ALERT_RED)

    # =========================================================================
    # SLIDE 2: Vision & Core Objectives
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_bg(s2)
    add_top_bar(s2, "PROJECT VISION & CORE OBJECTIVES")

    add_hud_card(s2, 0.8, 1.6, 5.7, 5.2, "SYSTEM VISION", [
        ("Iron Man JARVIS Concept", "Bringing a futuristic floating HUD assistant into a compact, affordable hardware device."),
        ("Physical-Digital Synergy", "Connecting natural language cloud AI directly to physical relay switches and hardware sensors."),
        ("Zero-Latency Interaction", "Sub-second speech processing for seamless human-computer interaction.")
    ], NEON_CYAN)

    add_hud_card(s2, 6.8, 1.6, 5.7, 5.2, "CORE PILLARS", [
        ("Edge Audio Capture", "16 kHz continuous sampling via ESP32 INMP441 mic streamed over raw TCP sockets."),
        ("Cloud AI Intelligence", "Groq Whisper STT + Llama 3.3 70B LLM with native function tool calling."),
        ("Optical Hologram HUD", "Pepper's Ghost 4-axis Arc Reactor graphic rendered on a 0.96\" OLED mini pyramid."),
        ("Smart IoT Control", "Optocoupled 5V relay module for switching real-world appliances (lights, fans, AC).")
    ], NEON_BLUE)

    # =========================================================================
    # SLIDE 3: 3-Tier System Architecture
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_bg(s3)
    add_top_bar(s3, "3-TIER SYSTEM ARCHITECTURE")

    add_hud_card(s3, 0.8, 1.6, 3.6, 5.2, "TIER 1: ESP32 EDGE NODE", [
        ("Microcontroller", "ESP32 WROOM-32 DevKit V1"),
        ("Audio Input", "INMP441 / Electret Mic (GPIO 34)"),
        ("Visual Output", "0.96\" SSD1306 OLED (I2C)"),
        ("Hardware Output", "5V Relay Module (GPIO 26)"),
        ("Network Protocol", "Wi-Fi TCP Socket Stream")
    ], NEON_CYAN)

    add_hud_card(s3, 4.8, 1.6, 3.6, 5.2, "TIER 2: HOST PROCESSING SERVER", [
        ("DSP Filter", "IIR DC-Blocker (R=0.995)"),
        ("VAD Engine", "Adaptive Noise Floor Tracker"),
        ("TCP Master", "Multi-Threaded Socket Server"),
        ("Audio Player", "Edge-TTS Pygame Lipsync Player"),
        ("HUD Controller", "State Synchronizer (TCP JSON)")
    ], NEON_BLUE)

    add_hud_card(s3, 8.8, 1.6, 3.7, 5.2, "TIER 3: GROQ CLOUD AI", [
        ("STT Engine", "Groq Whisper-Large-v3"),
        ("LLM Core", "Groq Llama-3.3-70b-Versatile"),
        ("Tool Dispatch", "JSON Function Calling (toggle_relay)"),
        ("Speech Synth", "Microsoft Edge Neural Voice"),
        ("Response Speed", "300+ Tokens/Sec Output")
    ], NEON_GREEN)

    # =========================================================================
    # SLIDE 4: Digital Signal Processing (DSP) & VAD
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_bg(s4)
    add_top_bar(s4, "DSP & ADAPTIVE VOICE ACTIVITY DETECTION (VAD)")

    add_hud_card(s4, 0.8, 1.6, 5.7, 5.2, "SIGNAL CONDITIONING (IIR DC-BLOCKER)", [
        ("ADC Bias Removal", "Filters out 2048 DC offset from ESP32 12-bit ADC input."),
        ("Filter Equation", "y[n] = x[n] - x[n-1] + R * y[n-1]  (Pole R = 0.995)."),
        ("Hum Suppression", "Cuts low-frequency electrical noise and mic pops without degrading speech pitch.")
    ], NEON_CYAN)

    add_hud_card(s4, 6.8, 1.6, 5.7, 5.2, "ADAPTIVE VAD ALGORITHM", [
        ("Noise Floor Tracking", "Smooths background noise using alpha coefficient 0.02."),
        ("Dynamic Threshold", "VAD Threshold = Noise Floor x 3.2 (prevents false triggers)."),
        ("Pre-Roll Buffer", "250ms circular audio buffer prepended to catch initial consonants."),
        ("Silence Hold", "900ms silence flush timer finishes user utterance gracefully.")
    ], NEON_BLUE)

    # =========================================================================
    # SLIDE 5: Cloud AI Engine & Native Tool Calling
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_bg(s5)
    add_top_bar(s5, "CLOUD AI REASONING & NATIVE TOOL CALLING")

    add_hud_card(s5, 0.8, 1.6, 5.7, 5.2, "GROQ CLOUD INTEGRATION", [
        ("Sub-Second STT", "Groq Whisper-Large-v3 transcribes speech in under 200ms."),
        ("High-Speed Intelligence", "Llama-3.3-70B provides instant, natural conversation."),
        ("Cost Efficiency", "Free tier API access provides professional cloud performance.")
    ], NEON_CYAN)

    add_hud_card(s5, 6.8, 1.6, 5.7, 5.2, "HARDWARE TOOL CALLING (FUNCTION DISPATCH)", [
        ("Natural Intent", "User says: 'Jarvis, turn on the fan.'"),
        ("LLM Tool Trigger", "Llama 3.3 auto-selects toggle_relay(device='fan', state='on')."),
        ("TCP Dispatch", "Server emits JSON payload {'command':'relay','state':'on'} to ESP32."),
        ("Visual Feedback", "HUD instantly switches to RELAY state with alert highlight.")
    ], ALERT_RED)

    # =========================================================================
    # SLIDE 6: IoT Hardware Control & Relay Circuit
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_bg(s6)
    add_top_bar(s6, "HARDWARE SWITCHING & RELAY MODULE")

    add_hud_card(s6, 0.8, 1.6, 5.7, 5.2, "5V RELAY HARDWARE ARCHITECTURE", [
        ("Optocoupler Isolation", "Protects ESP32 micro-controller from high-voltage AC spikes."),
        ("Control Signal", "Driven by ESP32 GPIO 26 pin."),
        ("Load Handling", "Switches up to 250V AC / 10A (Fans, Lamps, Appliances).")
    ], ALERT_RED)

    add_hud_card(s6, 6.8, 1.6, 5.7, 5.2, "COMMAND EXECUTION LOOP", [
        ("Step 1", "ESP32 captures voice request and streams PCM audio over TCP."),
        ("Step 2", "Python server executes Groq tool call and updates system state."),
        ("Step 3", "TCP socket sends JSON command back to ESP32 client."),
        ("Step 4", "ESP32 sets GPIO 26 HIGH, closing relay contacts and turning appliance ON.")
    ], NEON_BLUE)

    # =========================================================================
    # SLIDE 7: Pepper's Ghost Hologram Optics & Mini OLED
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_bg(s7)
    add_top_bar(s7, "PEPPER'S GHOST 4-AXIS OPTICAL HOLOGRAM")

    add_hud_card(s7, 0.8, 1.6, 5.7, 5.2, "OPTICAL PRINCIPLE (PEPPER'S GHOST)", [
        ("Partial Reflection", "Light projects upward from screen and reflects off 45-degree clear acrylic faces."),
        ("Black Background", "Pure black screen emits zero light, making glass transparent."),
        ("3D Spatial Perception", "Observer sees cyan Arc Reactor graphic floating in mid-air inside the pyramid.")
    ], NEON_CYAN)

    add_hud_card(s7, 6.8, 1.6, 5.7, 5.2, "STANDALONE MINI OLED IMPLEMENTATION", [
        ("Display Component", "0.96\" SSD1306 I2C OLED Screen (128x64 pixels)."),
        ("Reflector", "1-inch clear acrylic mini 4-sided pyramid sitting on OLED."),
        ("ESP32 Geometry", "Draws 4 mirrored Arc Reactor ring patterns in hardware via drawMiniHologramQuad().")
    ], NEON_GREEN)

    # =========================================================================
    # SLIDE 8: Neural Speech Synthesis & Audio Lipsync
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_bg(s8)
    add_top_bar(s8, "NEURAL VOICE SYNTHESIS & AMPLITUDE LIPSYNC")

    add_hud_card(s8, 0.8, 1.6, 5.7, 5.2, "MICROSOFT EDGE NEURAL TTS", [
        ("Voice Model", "en-US-ChristopherNeural (natural, authoritative assistant voice)."),
        ("Async Generation", "Synthesizes MP3 streams in background without blocking audio server."),
        ("Zero API Cost", "High-quality neural speech synthesis available free of charge.")
    ], NEON_BLUE)

    add_hud_card(s8, 6.8, 1.6, 5.7, 5.2, "AUDIO AMPLITUDE LIPSYNC", [
        ("Real-Time Modulation", "Audio energy level modulates the size and intensity of Arc Reactor core."),
        ("Visual States", "STANDBY (cyan pulse), THINKING (spin), SPEAKING (lipsync), RELAY (red flash)."),
        ("TCP Broadcast", "Server broadcasts state packets to ESP32 display loop.")
    ], NEON_CYAN)

    # =========================================================================
    # SLIDE 9: ESP32 Hardware Pin Wiring & Portal
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_bg(s9)
    add_top_bar(s9, "ESP32 PIN WIRING & WI-FI SETUP PORTAL")

    add_hud_card(s9, 0.8, 1.6, 5.7, 5.2, "HARDWARE PIN MAPPING", [
        ("GPIO 34 (ADC)", "Analog Microphone Audio Input (16 kHz sampling)."),
        ("GPIO 21 (SDA)", "0.96\" SSD1306 OLED Screen Data."),
        ("GPIO 22 (SCL)", "0.96\" SSD1306 OLED Screen Clock."),
        ("GPIO 26 (OUT)", "5V Relay Switch Control Signal."),
        ("5V & GND", "Common Power & Ground Rails.")
    ], NEON_CYAN)

    add_hud_card(s9, 6.8, 1.6, 5.7, 5.2, "EMBEDDED SETUP PORTAL", [
        ("AP Fallback Mode", "Starts AP 'Jarvis-Setup' at 192.168.4.1 if Wi-Fi disconnects."),
        ("NVS Storage", "Saves Wi-Fi SSID, password, server IP, and port in flash memory."),
        ("Instant Auto-Connect", "Reboots and connects to server automatically on power up.")
    ], NEON_GREEN)

    # =========================================================================
    # SLIDE 10: Performance Benchmarks & Highlights
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    apply_bg(s10)
    add_top_bar(s10, "PERFORMANCE BENCHMARKS & KEY HIGHLIGHTS")

    add_stat_badge(s10, 0.8, 1.6, 5.7, 2.4, "< 800 ms", "TOTAL VOICE-TO-ACTION LATENCY", NEON_CYAN)
    add_stat_badge(s10, 6.8, 1.6, 5.7, 2.4, "16.0 kHz", "TCP AUDIO STREAM FREQUENCY", NEON_BLUE)
    add_stat_badge(s10, 0.8, 4.4, 5.7, 2.4, "100%", "HANDS-FREE AUTOMATION", NEON_GREEN)
    add_stat_badge(s10, 6.8, 4.4, 5.7, 2.4, "< $15.00", "TOTAL HARDWARE BUILD COST", ALERT_RED)

    # =========================================================================
    # SLIDE 11: Future Roadmap & Upgrades
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    apply_bg(s11)
    add_top_bar(s11, "FUTURE ROADMAP & SYSTEM EXTENSIONS")

    add_hud_card(s11, 0.8, 1.6, 5.7, 5.2, "LOCAL OFFLINE AI (RASPBERRY PI 5)", [
        ("On-Device LLM", "Migrating from Groq cloud to local Ollama Qwen2.5 3B on Raspberry Pi 5."),
        ("Zero Internet Dependency", "Full voice recognition and hardware control operating 100% offline."),
        ("Privacy First", "Voice audio never leaves local home network.")
    ], NEON_BLUE)

    add_hud_card(s11, 6.8, 1.6, 5.7, 5.2, "HARDWARE EXTENSIONS", [
        ("MAX98357A I2S Audio Amp", "Direct speaker playback from ESP32 box."),
        ("Multi-Channel Relay Board", "Controlling 4 to 8 home automation channels."),
        ("Full-Body Hologram", "Scaling optical projection to 55-inch Pepper's Ghost floor display.")
    ], NEON_CYAN)

    # =========================================================================
    # SLIDE 12: Conclusion & Q&A
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    apply_bg(s12)

    frame12 = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    frame12.fill.solid()
    frame12.fill.fore_color.rgb = CARD_BG
    frame12.line.color.rgb = NEON_CYAN
    frame12.line.width = Pt(2)

    tb12 = s12.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "// PROJECT A.R.K. DEMONSTRATION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN
    p.font.name = "Consolas"
    p.space_after = Pt(14)

    p = tf12.add_paragraph()
    p.text = "THANK YOU!\nQUESTIONS & ANSWERS"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Arial"
    p.space_after = Pt(16)

    p = tf12.add_paragraph()
    p.text = "Project A.R.K. proves that futuristic holographic AI assistants with real-world IoT control can be built affordably using modern edge microcontrollers and cloud AI."
    p.font.size = Pt(14)
    p.font.color.rgb = NEON_BLUE
    p.font.name = "Arial"

    deck_path = "/home/kandan/Projects/mini_jarvis/PROJECT_ARK_Cyberpunk_Presentation.pptx"
    prs.save(deck_path)
    print("Presentation saved successfully at:", deck_path)

if __name__ == "__main__":
    create_cyberpunk_deck()
